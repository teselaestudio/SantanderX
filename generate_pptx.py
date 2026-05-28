import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use a blank slide layout (index 6 is usually blank in standard templates)
    blank_layout = prs.slide_layouts[6]

    # Colors
    bg_color = RGBColor(22, 22, 22)       # #161616
    text_color = RGBColor(224, 224, 224)  # #e0e0e0
    accent_color = RGBColor(16, 185, 129) # #10b981
    muted_color = RGBColor(136, 136, 136)  # #888888
    white_color = RGBColor(255, 255, 255)

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_title(slide, number_str, title_str, subtitle_str):
        # Number
        num_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(4), Inches(0.8))
        tf_num = num_box.text_frame
        tf_num.word_wrap = True
        p_num = tf_num.paragraphs[0]
        p_num.text = number_str
        p_num.font.name = 'Arial'
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = RGBColor(30, 41, 35) # Dark green tint for number background

        # Title and Subtitle
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(4.2), Inches(5.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = title_str
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = white_color
        p_title.space_after = Pt(10)

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle_str
        p_sub.font.name = 'Arial'
        p_sub.font.size = Pt(18)
        p_sub.font.color.rgb = accent_color
        p_sub.font.bold = True
        p_sub.space_after = Pt(20)
        
        return tf

    # ==========================================================================
    # SLIDE 1: PORTADA
    # ==========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)

    # Logo
    logo_path = 'Imagenes/Logo Tesela Estudio.png'
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(5.4), Inches(1.2), width=Inches(2.53))

    # Title & Subtitle Card
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.33), Inches(3.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "Cartografía digital para la gestión del territorio"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = 'Arial'
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = white_color
    p1.space_after = Pt(15)

    p2 = tf1.add_paragraph()
    p2.text = "Mapas colaborativos que conectan personas, empresas y administraciones para mejorar nuestras ciudades."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Arial'
    p2.font.size = Pt(20)
    p2.font.color.rgb = text_color
    p2.space_after = Pt(25)

    # Badges
    p3 = tf1.add_paragraph()
    p3.text = "ELEVATOR PITCH   •   SANTANDER X"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = 'Arial'
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = accent_color

    # Notes
    slide1.notes_slide.notes_text_frame.text = (
        "Presentación con logo de Tesela estudio.\n"
        "Título: “Cartografía digital para la gestión del territorio”\n"
        "Subtítulo: “Mapas colaborativos que conectan personas, empresas y administraciones para mejorar nuestras ciudades. “"
    )


    # ==========================================================================
    # SLIDE 2: ORIGEN - CONTEXTO
    # ==========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)

    tf2 = add_title(slide2, "01", "El Origen", "Investigación y Urbanismo")
    
    p_body = tf2.add_paragraph()
    p_body.text = "Como investigadores en el área de urbanismo y estudiantes de arquitectura, comprendimos la importancia de la relación entre el diseño urbano y la correcta gestión del territorio."
    p_body.font.name = 'Arial'
    p_body.font.size = Pt(16)
    p_body.font.color.rgb = text_color
    p_body.space_after = Pt(20)

    p_quote = tf2.add_paragraph()
    p_quote.text = "“Entender el espacio urbano es el primer paso para mejorarlo.”"
    p_quote.font.name = 'Arial'
    p_quote.font.size = Pt(16)
    p_quote.font.italic = True
    p_quote.font.color.rgb = white_color

    # Image
    img_path = 'Imagenes/Diapositiva 2.1/huerta murcia.jpg'
    if os.path.exists(img_path):
        slide2.shapes.add_picture(img_path, Inches(5.5), Inches(1.0), width=Inches(7.2), height=Inches(5.0))
        
        # Caption
        cap_box = slide2.shapes.add_textbox(Inches(5.5), Inches(6.1), Inches(7.2), Inches(0.5))
        cap_tf = cap_box.text_frame
        cap_p = cap_tf.paragraphs[0]
        cap_p.text = "Huerta de Murcia: Análisis territorial y urbanístico"
        cap_p.alignment = PP_ALIGN.CENTER
        cap_p.font.name = 'Arial'
        cap_p.font.size = Pt(11)
        cap_p.font.color.rgb = muted_color

    slide2.notes_slide.notes_text_frame.text = (
        "“1. El Origen\n"
        "Hace unos años, Alfonso y yo estábamos trabajando como investigadores en el área de urbanismo de la Universidad, "
        "durante la carrera de arquitectura, donde empezamos a trabajar en el diseño urbano y la gestión del territorio.\""
    )


    # ==========================================================================
    # SLIDE 3: ORIGEN - FRUSTRACIÓN
    # ==========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)

    tf3 = add_title(slide3, "02", "El Origen", "La Frustración Analógica")
    
    bullets = [
        "Toma de datos a mano: Procesos lentos con papel y lápiz, propensos al error.",
        "Información congelada: Datos en cajones o discos duros sin actualizar.",
        "Realidad dinámica: El territorio cambia constantemente y la cartografía quedaba obsoleta."
    ]
    for b in bullets:
        p_b = tf3.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = text_color
        p_b.space_after = Pt(10)

    p_dest = tf3.add_paragraph()
    p_dest.text = "De esta frustración, y con la intención de mejorar las cosas, nació Tesela Estudio."
    p_dest.font.name = 'Arial'
    p_dest.font.size = Pt(15)
    p_dest.font.bold = True
    p_dest.font.color.rgb = accent_color
    p_dest.space_before = Pt(15)

    # Images side by side
    img1 = 'Imagenes/Diapositiva 2.2/a mano 1.png'
    img2 = 'Imagenes/Diapositiva 2.2/a mano 2.png'
    
    if os.path.exists(img1):
        slide3.shapes.add_picture(img1, Inches(5.5), Inches(1.5), width=Inches(3.5), height=Inches(4.0))
    if os.path.exists(img2):
        slide3.shapes.add_picture(img2, Inches(9.2), Inches(1.5), width=Inches(3.5), height=Inches(4.0))
        
    cap_box = slide3.shapes.add_textbox(Inches(5.5), Inches(5.6), Inches(7.2), Inches(0.5))
    cap_p = cap_box.text_frame.paragraphs[0]
    cap_p.text = "Registro manual analógico y apuntes caóticos de campo"
    cap_p.alignment = PP_ALIGN.CENTER
    cap_p.font.name = 'Arial'
    cap_p.font.size = Pt(11)
    cap_p.font.color.rgb = muted_color

    slide3.notes_slide.notes_text_frame.text = (
        "\"Allí comenzamos a realizar las primeras cartografías tanto en el ámbito académico, como con empresas privadas "
        "dedicadas a la gestión de servicios urbanos. Al principio la toma de datos la hacíamos con papel y lápiz, casi de "
        "forma artesanal, lo que era un caos que acababa en un cajón o congelado en un disco duro mientras que la realidad "
        "que pretendía representar seguía en constante cambio, dejándonos la sensación agridulce de que al terminar, "
        "había que volver a empezar de nuevo. De esta frustración, y con la intención de mejorar las cosas, nació Tesela Estudio.\""
    )


    # ==========================================================================
    # SLIDE 4: PROBLEMA - PLATAFORMAS NUBE
    # ==========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)

    tf4 = add_title(slide4, "03", "El Problema", "Herramientas en la Nube")
    
    p_body = tf4.add_paragraph()
    p_body.text = "Descubrimos la existencia de plataformas colaborativas en la nube que centralizaban datos cartográficos y permitían conectar dispositivos."
    p_body.font.name = 'Arial'
    p_body.font.size = Pt(15)
    p_body.font.color.rgb = text_color
    p_body.space_after = Pt(20)

    # Competitor logos grid (3 cols x 2 rows)
    logos = [
        'Imagenes/Diapositiva 3.1/logo carto.png',
        'Imagenes/Diapositiva 3.1/logo giscloud.png',
        'Imagenes/Diapositiva 3.1/logo mapbox.png',
        'Imagenes/Diapositiva 3.1/logo merginmaps.png',
        'Imagenes/Diapositiva 3.1/logo movisat.png',
        'Imagenes/Diapositiva 3.1/logo qfield.png'
    ]
    
    grid_left = Inches(5.8)
    grid_top = Inches(1.5)
    col_width = Inches(2.2)
    row_height = Inches(1.8)
    
    for i, logo in enumerate(logos):
        if os.path.exists(logo):
            row = i // 3
            col = i % 3
            l = grid_left + col * (col_width + Inches(0.2))
            t = grid_top + row * (row_height + Inches(0.2))
            slide4.shapes.add_picture(logo, l, t, width=col_width)

    cap_box = slide4.shapes.add_textbox(Inches(5.5), Inches(5.5), Inches(7.2), Inches(0.5))
    cap_p = cap_box.text_frame.paragraphs[0]
    cap_p.text = "Plataformas de cartografía y SIG colaborativas analizadas"
    cap_p.alignment = PP_ALIGN.CENTER
    cap_p.font.name = 'Arial'
    cap_p.font.size = Pt(11)
    cap_p.font.color.rgb = muted_color

    slide4.notes_slide.notes_text_frame.text = (
        "“2. El problema\n"
        "Descubrimos que ya existían plataformas colaborativas en la nube que permitían conectar diferentes dispositivos "
        "a un mismo mapa. Por lo que ya no necesitábamos más pilas de papeles y toda la información la teníamos controlada.”"
    )


    # ==========================================================================
    # SLIDE 5: PROBLEMA - COMPLEJIDAD
    # ==========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)

    tf5 = add_title(slide5, "04", "El Problema", "Barreras de Usabilidad")
    
    p_body = tf5.add_paragraph()
    p_body.text = "El gran inconveniente detectado en el software existente en el mercado:"
    p_body.font.name = 'Arial'
    p_body.font.size = Pt(15)
    p_body.font.color.rgb = text_color
    p_body.space_after = Pt(15)

    p_quote = tf5.add_paragraph()
    p_quote.text = "“Las plataformas resultan poco intuitivas y complejas para usuarios no especializados en cartografía.”"
    p_quote.font.name = 'Arial'
    p_quote.font.size = Pt(16)
    p_quote.font.bold = True
    p_quote.font.italic = True
    p_quote.font.color.rgb = RGBColor(255, 82, 82)
    p_quote.space_after = Pt(20)

    bullets = [
        "Curva de aprendizaje empinada para operarios de campo.",
        "Pérdida de eficiencia al introducir datos en sistemas genéricos.",
        "Falta de adaptación al flujo real de los servicios urbanos."
    ]
    for b in bullets:
        p_b = tf5.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = text_color
        p_b.space_after = Pt(8)

    # Image
    img_path = 'Imagenes/Diapositiva 3.2/ejemplo movisat.png'
    if os.path.exists(img_path):
        slide5.shapes.add_picture(img_path, Inches(5.5), Inches(1.0), width=Inches(7.2), height=Inches(5.0))
        
        cap_box = slide5.shapes.add_textbox(Inches(5.5), Inches(6.1), Inches(7.2), Inches(0.5))
        cap_p = cap_box.text_frame.paragraphs[0]
        cap_p.text = "Ejemplo de interfaz técnica compleja de la competencia"
        cap_p.alignment = PP_ALIGN.CENTER
        cap_p.font.name = 'Arial'
        cap_p.font.size = Pt(11)
        cap_p.font.color.rgb = muted_color

    slide5.notes_slide.notes_text_frame.text = (
        "“Sin embargo, estas plataformas siempre resultaban poco intuitivas y complicadas de manejar para las personas "
        "que no estaban dentro del mundo de la cartografía y el territorio.”"
    )


    # ==========================================================================
    # SLIDE 6: PROBLEMA - SOLUCIÓN INICIAL (PILOTO ARBOLADO)
    # ==========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)

    tf6 = add_title(slide6, "05", "El Problema", "Piloto: Gestión de Arbolado")
    
    bullets = [
        "Interfaz amigable: Diseñada específicamente para trabajadores de campo.",
        "Registro de campo directo: Anotaciones y modificaciones sobre el terreno en tiempo real.",
        "Canal directo: Conexión inmediata oficina-campo, optimizando el flujo de información."
    ]
    for b in bullets:
        p_b = tf6.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = text_color
        p_b.space_after = Pt(15)

    p_accent = tf6.add_paragraph()
    p_accent.text = "Nosotros realizamos la cartografía y el sistema de gestión personalizado."
    p_accent.font.name = 'Arial'
    p_accent.font.size = Pt(15)
    p_accent.font.bold = True
    p_accent.font.color.rgb = accent_color

    # Image
    img_path = 'Imagenes/Diapositiva 3.3/ejemplo cartagena.png'
    if os.path.exists(img_path):
        slide6.shapes.add_picture(img_path, Inches(5.5), Inches(1.0), width=Inches(7.2), height=Inches(5.0))
        
        cap_box = slide6.shapes.add_textbox(Inches(5.5), Inches(6.1), Inches(7.2), Inches(0.5))
        cap_p = cap_box.text_frame.paragraphs[0]
        cap_p.text = "Plataforma piloto implementada en Cartagena"
        cap_p.alignment = PP_ALIGN.CENTER
        cap_p.font.name = 'Arial'
        cap_p.font.size = Pt(11)
        cap_p.font.color.rgb = muted_color

    slide6.notes_slide.notes_text_frame.text = (
        "“Aprovechando que ya estábamos colaborando con esta empresa nos propusimos utilizar estas plataformas para "
        "proponerles un sistema de gestión de arbolado accesible y legible a sus trabajadores .En el que cada operario "
        "pudiera anotar su trabajo y las modificaciones que realizaba. Comunicándose directamente con la oficina y el "
        "resto del personal. Nosotros realizamos la cartografía y el sistema de gestión personalizado. A ellos les facilitamos el flujo de información.”"
    )


    # ==========================================================================
    # SLIDE 7: LA SOLUCIÓN - PRODUCTO
    # ==========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)

    tf7 = add_title(slide7, "06", "La Solución", "Nuestro Producto")
    
    p_body = tf7.add_paragraph()
    p_body.text = "Reducimos la multiplicidad de documentos en distintos terminales y los errores de coordinación."
    p_body.font.name = 'Arial'
    p_body.font.size = Pt(15)
    p_body.font.color.rgb = text_color
    p_body.space_after = Pt(20)

    p_adv = tf7.add_paragraph()
    p_adv.text = "Nuestra Ventaja Competitiva:"
    p_adv.font.name = 'Arial'
    p_adv.font.size = Pt(16)
    p_adv.font.bold = True
    p_adv.font.color.rgb = accent_color
    p_adv.space_after = Pt(8)

    p_adv_desc = tf7.add_paragraph()
    p_adv_desc.text = "Conocemos a fondo el funcionamiento del urbanismo y las ciudades, lo que nos permite desarrollar sistemas adaptados a las necesidades reales del servicio."
    p_adv_desc.font.name = 'Arial'
    p_adv_desc.font.size = Pt(14)
    p_adv_desc.font.color.rgb = text_color

    # Images side by side (Phone & Desktop Visor)
    img_phone = 'Imagenes/Diapositiva 4.1/mockup telefono.png'
    img_visor = 'Imagenes/Diapositiva 4.1/visor municipio.png'
    
    if os.path.exists(img_phone):
        slide7.shapes.add_picture(img_phone, Inches(5.5), Inches(1.0), width=Inches(2.4), height=Inches(4.8))
    if os.path.exists(img_visor):
        slide7.shapes.add_picture(img_visor, Inches(8.1), Inches(1.5), width=Inches(4.6), height=Inches(3.5))

    cap_box = slide7.shapes.add_textbox(Inches(5.5), Inches(5.9), Inches(7.2), Inches(0.5))
    cap_p = cap_box.text_frame.paragraphs[0]
    cap_p.text = "App móvil del operario y visor de control de oficina"
    cap_p.alignment = PP_ALIGN.CENTER
    cap_p.font.name = 'Arial'
    cap_p.font.size = Pt(11)
    cap_p.font.color.rgb = muted_color

    slide7.notes_slide.notes_text_frame.text = (
        "“3. La solución\n"
        "De esta manera redujimos la multiplicidad de documentos en distintos terminales y los errores de coordinación que "
        "acarrea este tipo de funcionamiento. Nuestra ventaja competitiva a la hora de realizar las bases de datos es que "
        "conocemos el funcionamiento del urbanismo y las ciudades, lo que nos permite desarrollar un sistema consecuente "
        "para la gestión de sus servicios. Que es la base de nuestro producto.”"
    )


    # ==========================================================================
    # SLIDE 8: LA SOLUCIÓN - CASOS DE ÉXITO
    # ==========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8)

    tf8 = add_title(slide8, "07", "La Solución", "Sistemas en Marcha")
    
    p_body = tf8.add_paragraph()
    p_body.text = "Conseguimos hacer legibles los sistemas de la ciudad para que cualquier operario pueda actualizar atributos en tiempo real."
    p_body.font.name = 'Arial'
    p_body.font.size = Pt(15)
    p_body.font.color.rgb = text_color
    p_body.space_after = Pt(20)

    p_sub = tf8.add_paragraph()
    p_sub.text = "Tres servicios municipales actualmente activos y funcionando en el territorio:"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = muted_color
    p_sub.space_after = Pt(10)

    # 3 Cases side-by-side
    cases = [
        ('Imagenes/Diapositiva 4.2/Salobrena.jpg', 'Limpieza Viaria', Inches(5.3)),
        ('Imagenes/Diapositiva 4.2/funes.jpg', 'Arbolado Urbano', Inches(7.95)),
        ('Imagenes/Diapositiva 4.2/rosa.jpg', 'Recogida de Residuos', Inches(10.6))
    ]
    
    for img_path, label, left_pos in cases:
        if os.path.exists(img_path):
            slide8.shapes.add_picture(img_path, left_pos, Inches(1.8), width=Inches(2.5), height=Inches(3.33))
            
            # Label
            lbl_box = slide8.shapes.add_textbox(left_pos, Inches(5.2), Inches(2.5), Inches(0.5))
            lbl_tf = lbl_box.text_frame
            lbl_p = lbl_tf.paragraphs[0]
            lbl_p.text = label
            lbl_p.alignment = PP_ALIGN.CENTER
            lbl_p.font.name = 'Arial'
            lbl_p.font.size = Pt(13)
            lbl_p.font.bold = True
            lbl_p.font.color.rgb = accent_color

    slide8.notes_slide.notes_text_frame.text = (
        "\"Con nuestros mapas conseguimos hacer legible los sistemas que componen la ciudad o el territorio para que "
        "cualquier persona, desde un técnico a un trabajador de campo pueda editar atributos y tomar decisiones estratégicas "
        "en tiempo real. Creando una interfaz de uso específica para cada caso de gestión, como el arbolado, la limpieza "
        "viaria o la recogida de residuos. Que son ejemplos que ya tenemos en marcha y funcionando.\""
    )


    # ==========================================================================
    # SLIDE 9: ESCALABILIDAD
    # ==========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9)

    tf9 = add_title(slide9, "08", "Escalabilidad", "Negocio y Crecimiento")
    
    bullets = [
        "SaaS (Suscripción): Modelo de negocio escalable con ingresos recurrentes.",
        "Costes optimizados: Estructura de costes ágil sin barreras físicas.",
        "Expansión en desarrollo: Núcleo activo en Cartagena y nuevos municipios en fase de implantación.",
        "Clientes: Válido para administraciones públicas y empresas contratistas de servicios urbanos."
    ]
    for b in bullets:
        p_b = tf9.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = text_color
        p_b.space_after = Pt(12)

    # Map Image
    img_path = 'Imagenes/Diapositiva 5/spain.png'
    if os.path.exists(img_path):
        slide9.shapes.add_picture(img_path, Inches(5.5), Inches(1.0), width=Inches(7.2), height=Inches(5.0))
        
        cap_box = slide9.shapes.add_textbox(Inches(5.5), Inches(6.1), Inches(7.2), Inches(0.5))
        cap_p = cap_box.text_frame.paragraphs[0]
        cap_p.text = "municipios cartografiados"
        cap_p.alignment = PP_ALIGN.CENTER
        cap_p.font.name = 'Arial'
        cap_p.font.size = Pt(12)
        cap_p.font.bold = True
        cap_p.font.color.rgb = text_color

    slide9.notes_slide.notes_text_frame.text = (
        "“4. Escalabilidad\n"
        "El más importante está en Cartagena, pero tenemos otros municipios de la región en fase de desarrollo. "
        "Con planes de implementación posterior a mayor escala. Nuestro producto se basa en un modelo de suscripción, "
        "por lo que nuestra estructura de costes está optimizada y no depende de barreras físicas. Estos servicios no son "
        "solo válidos para este tipo de empresas, sino que también puede ser conveniente para las administraciones y "
        "cualquier tipo de entidad que trabaje con la ciudad y el territorio. Gracias a la estructura digital somos capaces "
        "de expandirnos por diferentes localizaciones de forma rápida y ágil.”"
    )


    # ==========================================================================
    # SLIDE 10: CIERRE CENTRADO
    # ==========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10)

    # Centered Logo
    if os.path.exists(logo_path):
        slide10.shapes.add_picture(logo_path, Inches(5.4), Inches(0.8), width=Inches(2.53))

    # Quote box in center
    quote_box = slide10.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(4.5))
    tf10 = quote_box.text_frame
    tf10.word_wrap = True
    
    p1 = tf10.paragraphs[0]
    p1.text = "“Lo que empezó como un problema para la toma de datos a papel y lápiz, lo hemos transformado en una herramienta capaz de coordinar personas, empresas y administraciones públicas para mejorar la gestión de nuestras ciudades y el territorio natural.”"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = 'Arial'
    p1.font.size = Pt(20)
    p1.font.italic = True
    p1.font.color.rgb = white_color
    p1.space_after = Pt(20)

    p2 = tf10.add_paragraph()
    p2.text = "Somos Tesela Estudio y gracias por escucharnos."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = accent_color
    p2.space_after = Pt(30)

    # Contact Info
    p3 = tf10.add_paragraph()
    p3.text = "WEB: teselaestudio.es     |     EMAIL: info.tesela.estudio@gmail.com"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = text_color

    slide10.notes_slide.notes_text_frame.text = (
        "\"Lo que empezó como un problema nuestro para la toma de datos al empezar con papel y lápiz, lo hemos "
        "transformado en una herramienta capaz de coordinar a personas, empresas y administraciones públicas para mejorar "
        "la gestión de nuestras ciudades y el territorio natural. Somos Tesela Estudio y gracias por escucharnos\""
    )

    # Save presentation
    output_path = 'Presentacion_Tesela_Estudio.pptx'
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == '__main__':
    create_presentation()
